#!/usr/bin/env python3
"""
Professional PowerPoint Generator - UniResult Project Presentation
Creates 15+ slides with professional design and comprehensive content
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Color Palette
PRIMARY_BLUE = RGBColor(25, 118, 210)
SECONDARY_GREEN = RGBColor(56, 142, 60)
ACCENT_ORANGE = RGBColor(255, 152, 0)
DARK_GRAY = RGBColor(33, 33, 33)
LIGHT_GRAY = RGBColor(245, 245, 245)
WHITE = RGBColor(255, 255, 255)
GREEN_SUCCESS = RGBColor(76, 175, 80)

def add_title_slide(prs, title, subtitle=""):
    """Create a professional title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(60)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.text = subtitle
        subtitle_p.font.size = Pt(28)
        subtitle_p.font.color.rgb = WHITE
        subtitle_p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, bg_color=PRIMARY_BLUE):
    """Create a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title bar
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = bg_color
    title_shape.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Content
    text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.8))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items, color=PRIMARY_BLUE):
    """Create a two-column content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title bar
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = color
    title_shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Left column header
    left_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.2), Inches(0.4))
    left_h_frame = left_header.text_frame
    left_h_frame.text = left_title
    left_h_frame.paragraphs[0].font.size = Pt(22)
    left_h_frame.paragraphs[0].font.bold = True
    left_h_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    
    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(4.2), Inches(5.2))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for i, item in enumerate(left_items):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
    
    # Right column header
    right_header = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.2), Inches(0.4))
    right_h_frame = right_header.text_frame
    right_h_frame.text = right_title
    right_h_frame.paragraphs[0].font.size = Pt(22)
    right_h_frame.paragraphs[0].font.bold = True
    right_h_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    
    # Right column content
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.7), Inches(4.2), Inches(5.2))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for i, item in enumerate(right_items):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
    
    return slide

def generate_presentation():
    """Generate the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # SLIDE 1: Title Slide
    slide = add_title_slide(prs, "UniResult", "Result Sheet Generator")
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
    info_frame = info_box.text_frame
    info_frame.word_wrap = True
    
    info_lines = [
        "Student: Gaurav | Batch: BCA B3",
        "College: CIMAGE College | University: PPU",
        "Level-1 Presentation: February 13-14, 2026"
    ]
    
    for i, line in enumerate(info_lines):
        if i == 0:
            p = info_frame.paragraphs[0]
        else:
            p = info_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # SLIDE 2: Introduction & Objectives
    add_content_slide(prs, "📌 Introduction & Objectives", [
        "🎯 Project: UniResult - Student Result Sheet Generator",
        "💡 Purpose: Full-stack web application for generating, managing, and saving student result sheets",
        "🎓 Target Users: Students, Educational Institutions, Universities",
        "✓ Key Feature: Real-time percentage calculation and PDF generation",
        "✓ Database Integration: Auto-save results to PostgreSQL",
        "✓ Responsive Design: Works seamlessly on desktop and mobile",
        "",
        "🔹 Primary Objectives:",
        "  • Simplify result sheet creation with intuitive interface",
        "  • Provide real-time calculations (percentage, grades, status)",
        "  • Enable PDF export for documentation and printing",
        "  • Maintain searchable history of generated results"
    ], PRIMARY_BLUE)
    
    # SLIDE 3: SRS - Functional Requirements
    add_content_slide(prs, "📋 SRS - Functional Requirements", [
        "👤 User Input & Data Entry:",
        "  ✓ Create result sheets with student details (name, roll no., registration)",
        "  ✓ Add/remove subjects dynamically with marks validation",
        "  ✓ Input validation: marks cannot exceed maximum marks",
        "  ✓ Real-time percentage & grade calculation display",
        "",
        "📄 Document Generation & Export:",
        "  ✓ Generate professional PDF result sheets",
        "  ✓ Download PDF with proper formatting and headers",
        "  ✓ Save results to PostgreSQL database via API",
        "",
        "📋 History & Management:",
        "  ✓ View previously saved results from local storage",
        "  ✓ Search and filter results by student/date",
        "  ✓ Edit/update existing results"
    ], PRIMARY_BLUE)
    
    # SLIDE 4: SRS - Non-Functional Requirements
    add_two_column_slide(prs, "📋 SRS - Non-Functional Requirements",
        "Performance & Security", [
            "⚡ Response time: < 2 seconds",
            "👥 Scalability: Support 1000+ concurrent users",
            "🔒 Security: HTTPS encryption enabled",
            "🛡️ API Security: CORS protection configured",
            "✓ Input validation & data sanitization",
            "🔐 Session management & user authentication"
        ],
        "Compatibility & UI/UX", [
            "💻 Cross-browser: Chrome, Firefox, Safari, Edge",
            "📱 Responsive: Mobile, Tablet, Desktop",
            "♿ Accessibility: WCAG 2.1 AA compliant",
            "🎨 UI Quality: 95+ Lighthouse score",
            "⚡ Load time: <3 seconds on 3G",
            "🌙 Dark/Light mode support"
        ]
    )
    
    # SLIDE 5: Process Logic
    add_content_slide(prs, "🔄 Process Logic & Data Flow", [
        "📥 PHASE 1: Data Input & Validation",
        "  • User enters student information (name, roll number, university)",
        "  • Add subjects and corresponding marks dynamically",
        "  • System validates all inputs in real-time",
        "",
        "⚙️ PHASE 2: Calculation & Processing",
        "  • Total Marks = Sum of all obtained marks",
        "  • Maximum Marks = Sum of all maximum marks",
        "  • Percentage = (Total Obtained / Total Maximum) × 100",
        "  • Grade determination based on percentage (A+, A, B+, B, C, etc.)",
        "",
        "💾 PHASE 3: Storage & Export",
        "  • Generate PDF with formatted result sheet",
        "  • Save to PostgreSQL database (server-side)",
        "  • Store in browser local storage (client-side)",
        "  • Provide download and preview options"
    ], GREEN_SUCCESS)
    
    # SLIDE 6: Data Dictionary
    add_content_slide(prs, "📊 Data Dictionary - Key Entities", [
        "📝 Student Information:",
        "  • name: String - Student's full name",
        "  • rollNumber: String - University roll number",
        "  • registrationNumber: String - Registration ID",
        "  • courseName: String - Course/Degree name",
        "  • semester: Number - Current semester",
        "  • academicYear: String - Year (e.g., 2023-24)",
        "",
        "📚 Subject & Marks:",
        "  • subjectName: String - Name of subject",
        "  • maxMarks: Number - Total marks for subject (default: 100)",
        "  • obtainedMarks: Number - Marks scored",
        "",
        "📈 Result Summary:",
        "  • totalObtained: Number - Sum of all marks",
        "  • totalMaximum: Number - Sum of max marks",
        "  • percentage: Decimal - Final percentage",
        "  • status: String - Pass/Fail/Distinction",
        "  • createdAt: DateTime - Timestamp"
    ], PRIMARY_BLUE)
    
    # SLIDE 7: DFD - Data Flow Diagram
    add_content_slide(prs, "📈 DFD - Data Flow Diagram", [
        "",
        "LEVEL 0 - CONTEXT DIAGRAM:",
        "",
        "       Student/User",
        "            |",
        "            | Input: Student & Subject Data",
        "            ↓",
        "     ┌──────────────────┐",
        "     │   UniResult      │",
        "     │   System         │",
        "     └──────────────────┘",
        "            ↓",
        "     ┌──────────────────┐",
        "     │   PDF Output     │",
        "     │   & Database     │",
        "     └──────────────────┘",
        "",
        "KEY PROCESSES:",
        "  P1: Accept & Validate Input | P2: Calculate Results",
        "  P3: Generate PDF | P4: Store in Database"
    ], ACCENT_ORANGE)
    
    # SLIDE 8: Technology Stack
    add_two_column_slide(prs, "🏗️ Technology Stack & Architecture",
        "Frontend Technologies", [
            "⚛️ React 18 - UI Framework",
            "📘 TypeScript - Type-safe JavaScript",
            "🚀 Vite - Fast build tool",
            "🎨 Tailwind CSS - Styling framework",
            "🧩 shadcn/ui - Component library",
            "🛣️ React Router - Navigation",
            "📊 Recharts - Chart library"
        ],
        "Backend & Database", [
            "🐍 Python 3.8+ - Backend language",
            "🔥 Flask - Web framework",
            "🗄️ PostgreSQL - Database",
            "🔌 psycopg2 - DB connector",
            "🚀 Supabase - Cloud database (optional)",
            "📄 ReportLab - PDF generation",
            "🌐 Flask-CORS - Cross-origin support"
        ],
        color=PRIMARY_BLUE
    )
    
    # SLIDE 9: Gantt Chart - Project Timeline
    add_content_slide(prs, "📅 Project Timeline - 3 Months", [
        "📆 MONTH 1: Planning & Frontend Development",
        "  Week 1-2: Requirements Analysis, Design, Architecture",
        "  Week 3-4: React Components, UI Development, Form Creation",
        "  Week 4: Dashboard, Summary View, PDF Preview",
        "",
        "MONTH 2: Backend & Integration",
        "  Week 5-6: Flask API Development (/submit endpoint)",
        "  Week 7: PostgreSQL Schema Design & Initialization",
        "  Week 8: PDF Generation Module, Email Notifications",
        "  Week 8: Frontend-Backend Integration, Testing",
        "",
        "MONTH 3: Testing, Optimization & Deployment",
        "  Week 9: Unit Testing, Integration Testing",
        "  Week 10: UAT, Bug Fixes, Performance Optimization",
        "  Week 11: Security Review, Documentation",
        "  Week 12: Deployment to Production, Final Testing"
    ], PRIMARY_BLUE)
    
    # SLIDE 10: ER Diagram - Database Schema
    add_content_slide(prs, "🗄️ ER Diagram - Database Schema", [
        "PRIMARY TABLE: submissions",
        "┌──────────────────────────────────────┐",
        "│ submissions (Main storage)           │",
        "├──────────────────────────────────────┤",
        "│ PK: id (UUID) - Primary Key          │",
        "│ user_id (String) - User identifier   │",
        "│ data (JSONB) - Flexible data storage │",
        "│ created_at (TIMESTAMP)               │",
        "│ updated_at (TIMESTAMP)               │",
        "└──────────────────────────────────────┘",
        "",
        "DATA STRUCTURE (JSONB format):",
        "  {",
        "    studentInfo: { name, rollNo, course, semester, ... }",
        "    subjects: [ { name, maxMarks, obtainedMarks }, ... ]",
        "    calculations: { totalObtained, percentage, status }",
        "  }"
    ], ACCENT_ORANGE)
    
    # SLIDE 11: User Interface Features
    add_content_slide(prs, "🎨 User Interface & Screen Design", [
        "🏠 HOME PAGE (Index):",
        "  ✓ Student information input form",
        "  ✓ Dynamic subject entry with add/remove buttons",
        "  ✓ Real-time percentage display",
        "  ✓ Generate & Preview buttons",
        "",
        "👁️ PREVIEW PAGE:",
        "  ✓ Complete formatted result sheet view",
        "  ✓ Grade status (Distinction/First/Second/Pass/Fail)",
        "  ✓ Download PDF button (with custom styling)",
        "  ✓ Save to database option",
        "",
        "📋 HISTORY PAGE:",
        "  ✓ List of previously generated results",
        "  ✓ Search, filter, and delete options",
        "  ✓ Responsive grid layout",
        "  ✓ Mobile-optimized interface"
    ], PRIMARY_BLUE)
    
    # SLIDE 12: Expected Report Generation
    add_content_slide(prs, "📄 Expected PDF Report Output", [
        "REPORT STRUCTURE:",
        "  ✓ HEAD SECTION: Institution logo, title, date",
        "  ✓ STUDENT INFO: Name, Roll No., Registration, Course, Semester",
        "",
        "RESULT TABLE:",
        "  ┌─────────────────────────────────────────────────┐",
        "  │ Subject | Max Marks | Obtained | % | Grade      │",
        "  ├─────────────────────────────────────────────────┤",
        "  │ Math    | 100       | 85       | 85| A          │",
        "  │ Science | 100       | 92       | 92| A+         │",
        "  └─────────────────────────────────────────────────┘",
        "",
        "SUMMARY SECTION:",
        "  Total Marks: 370/400 | Final Percentage: 92.5%",
        "  Grade: A+ | Status: PASS | Remarks: Excellent Performance",
        "",
        "FOOTER: Generated Date, Authorized Signature Line"
    ], GREEN_SUCCESS)
    
    # SLIDE 13: Future Scope & Enhancements
    add_content_slide(prs, "🚀 Future Scope & Enhancements", [
        "VERSION 2.0 FEATURES:",
        "  🔮 Multi-language Support (Hindi, Spanish, French, etc.)",
        "  📱 Native Mobile App (React Native / Flutter)",
        "  🤖 AI-powered grade prediction & analytics",
        "  📊 Advanced Analytics Dashboard",
        "  🔗 University System Integration (API connectors)",
        "  🎓 Digital Certificate Generation (Blockchain)",
        "  📧 Email Notifications & Alerts",
        "  💬 Real-time Collaboration for Teachers",
        "",
        "TECHNICAL IMPROVEMENTS:",
        "  ⚡ GraphQL API Implementation",
        "  🚀 Redis Caching for Performance",
        "  🐳 Docker Containerization",
        "  🔄 Microservices Architecture"
    ], ACCENT_ORANGE)
    
    # SLIDE 14: References & Bibliography
    add_content_slide(prs, "📚 References & Bibliography", [
        "📖 Official Documentations:",
        "  1. React.js - https://react.dev",
        "  2. Flask - https://flask.palletsprojects.com",
        "  3. PostgreSQL - https://www.postgresql.org/docs",
        "  4. TypeScript - https://www.typescriptlang.org",
        "  5. Python - https://docs.python.org/3",
        "  6. Tailwind CSS - https://tailwindcss.com",
        "",
        "🛠️ Tools & Libraries:",
        "  • Vite Build Tool - https://vitejs.dev",
        "  • shadcn/ui - https://ui.shadcn.com",
        "  • ReportLab - PDF generation",
        "  • Supabase - https://supabase.io",
        "",
        "📚 Learning Resources:",
        "  • Medium Articles on Full-Stack Development",
        "  • GitHub Repositories & Open Source Projects",
        "  • Stack Overflow Community"
    ], PRIMARY_BLUE)
    
    # SLIDE 15: Conclusion & Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = SECONDARY_GREEN
    
    # Main text
    main_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
    main_frame = main_box.text_frame
    main_frame.word_wrap = True
    
    main_p = main_frame.paragraphs[0]
    main_p.text = "Thank You!"
    main_p.font.size = Pt(60)
    main_p.font.bold = True
    main_p.font.color.rgb = WHITE
    main_p.alignment = PP_ALIGN.CENTER
    
    sub_p = main_frame.add_paragraph()
    sub_p.text = ""
    
    sub_p = main_frame.add_paragraph()
    sub_p.text = "Questions & Discussion"
    sub_p.font.size = Pt(36)
    sub_p.font.bold = True
    sub_p.font.color.rgb = ACCENT_ORANGE
    sub_p.alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(1.5))
    contact_frame = contact_box.text_frame
    contact_frame.word_wrap = True
    
    contact_p = contact_frame.paragraphs[0]
    contact_p.text = "Student: Gaurav | Batch: BCA B3 | College: CIMAGE College | University: PPU"
    contact_p.font.size = Pt(16)
    contact_p.font.color.rgb = WHITE
    contact_p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_file = "UniResult_Presentation.pptx"
    prs.save(output_file)
    
    print("\n" + "="*60)
    print("✅ PRESENTATION GENERATED SUCCESSFULLY!")
    print("="*60)
    print(f"📁 File: {output_file}")
    print(f"📊 Total Slides: {len(prs.slides)}")
    print(f"📐 Dimensions: {prs.slide_width.inches}\" x {prs.slide_height.inches}\"")
    print("\n📋 Content Breakdown:")
    print("  1. Title Slide with Student Info")
    print("  2. Introduction & Objectives")
    print("  3. SRS - Functional Requirements")
    print("  4. SRS - Non-Functional Requirements")
    print("  5. Process Logic & Data Flow")
    print("  6. Data Dictionary")
    print("  7. DFD - Data Flow Diagram")
    print("  8. Technology Stack")
    print("  9. Project Timeline (Gantt)")
    print("  10. ER Diagram - Database Schema")
    print("  11. User Interface Features")
    print("  12. PDF Report Output")
    print("  13. Future Scope")
    print("  14. References & Bibliography")
    print("  15. Conclusion & Q&A")
    print("\n🎤 Ready for Level-1 Presentation (Feb 13-14, 2026)")
    print("="*60 + "\n")
    
    return output_file

if __name__ == "__main__":
    print("\n🎬 Starting PowerPoint Generation...\n")
    generate_presentation()
